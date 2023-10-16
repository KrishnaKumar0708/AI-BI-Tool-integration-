from cProfile import label
import numpy as np
import requests
import json
import matplotlib as mpl
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR
from numpy import shape
from pptx.util import Pt
import os

# type: ignore
# pyright: reportMissingModuleSource=false
# pyright: reportMissingImports = false

repo_rev, repo_year, grs_prof, ebit, net_inc= [], [], [], [], []

prs = Presentation()

#################################################################################################################################
#################################################################################################################################

def Revenue_1():
    
    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)

    n = len(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'])

    '''with open("output.json", "w") as f:
        json.dump(response.json()['incomeStatementHistoryQuarterly']['totalRevenue']['fmt'][:-1], f)
        #json.dump(response.json()['majorHoldersBreakdown']['institutionsPercentHeld'], f)'''

    for i in range(n-1):
        repo_rev.append(float(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'][i]['totalRevenue']['fmt'][:-1]))
        repo_year.append('Q'+str(i+1)+"'"+str(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'][i]['endDate']['fmt'][2:4]))
    
    rev_per = []
    for i in range(1,n-1):
        rev_per.append(str((repo_rev[i]/repo_rev[i-1] - 1)*100)+'%')

    print(repo_year)
    print(repo_rev)
    print(rev_per)


    fig, ax = plt.subplots()
    width = 0.3
    ind = np.arange(len(repo_rev))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, repo_rev, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(repo_rev):
        ax.text(i-0.15,v+5, str(v))
    
    plt.plot(ind, repo_rev)
    for x,y in enumerate(repo_rev):
        ax.arrow(x,y+20,0.2,0.2,head_width = 0.02,width = 0.5)

    ax.set_xticks(ind, repo_year)
    #ax.set_title('Revenue                                                                                  ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Rev_1.png")

def Gross_prof():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    n = len(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'])

    for i in range(n-1) :
        grs_prof.append(float(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'][i]['grossProfit']['fmt'][:-1]))

    gross_per = []
    for i in range(1,n-1):
        gross_per.append(str((grs_prof[i]/grs_prof[i-1] - 1)*100)+'%')

    print(grs_prof)
    print(gross_per)


    fig, ax = plt.subplots()
    width = 0.3
    ind = np.arange(len(grs_prof))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, grs_prof, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(grs_prof):
        ax.text(i-0.15,v+2, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('Revenue                                                                                  ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("grs_profit_1.png")

def EBIT():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    ebit_per = []

    n = len(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'])

    for i in range(n-1) :
        ebit.append(float(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'][i]['ebit']['fmt'][:-1]))

    ebit_per = []
    for i in range(1,n-1) :
        ebit_per.append(str((ebit[i]/ebit[i-1] - 1)*100)+'%')

    print(ebit)
    print(ebit_per)

    fig, ax = plt.subplots()
    width = 0.3
    ind = np.arange(len(ebit))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ebit, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ebit):
        ax.text(i-0.15,v+1.5, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('Revenue                                                                                  ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Ebit_1.png")


def Net_Income():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    net_inc_per = []

    n = len(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'])

    for i in range(n-1) :
        net_inc.append(float(response.json()['incomeStatementHistoryQuarterly']['incomeStatementHistory'][i]['netIncomeApplicableToCommonShares']['fmt'][:-1]))
    
    net_inc_per = []
    for i in range(1,n-1) :
        net_inc_per.append(str((net_inc[i]/net_inc[i-1] - 1)*100)+'%')

    print(net_inc)
    print(net_inc_per)

    fig, ax = plt.subplots()
    width = 0.3
    ind = np.arange(len(net_inc))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.1)

    ax.bar(ind, net_inc, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(net_inc):
        ax.text(i-0.15,v+1, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('Revenue                                                                                  ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Net_inc_1.png")


def create_pptx_report4():
    
    Revenue_1()
    Gross_prof()
    EBIT()
    Net_Income()

    #prs = Presentation()
    blank_slide_layout1 = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout1)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    connector = slide1.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.1), Inches(1.9), Inches(1.1))
    color_line = connector.line
    #color_line.solid()
    color_line.color.rgb = RGBColor(250, 98, 28)
    
    #Text box for Marginal Outputs
    txBox = slide1.shapes.add_textbox(left=Inches(0.64), top=Inches(0.41), width=Inches(11.5), height=Inches(0.6))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Earnings'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(32)
    font.bold = False
    font.italic = None 
    font.color.rgb = RGBColor(19, 46, 87)

    pic = slide1.shapes.add_picture('Rev_1.png', Inches(1.5), Inches(1.6), width=Inches(6), height=Inches(3.5))
    pic = slide1.shapes.add_picture('grs_profit_1.png', Inches(8.5), Inches(1.6), width=Inches(6), height=Inches(3.5))
    pic = slide1.shapes.add_picture('Ebit_1.png', Inches(1.5), Inches(5.3), width=Inches(6), height=Inches(3.5))
    pic = slide1.shapes.add_picture('Net_inc_1.png', Inches(8.5), Inches(5.3), width=Inches(6), height=Inches(3.5))
    #pic = slide1.shapes.add_picture('FCF.png', Inches(8.5), Inches(5.3), width=Inches(6), height=Inches(3.5))

    #Text box for Revenue
    txBox = slide1.shapes.add_textbox(left=Inches(1.5), top=Inches(1.6), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Revenue'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)
    

    #Text box for Gross Profit
    txBox = slide1.shapes.add_textbox(left=Inches(8.5), top=Inches(1.6), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Gross Profit'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    #Text box for EBIT
    txBox = slide1.shapes.add_textbox(left=Inches(1.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EBIT'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    #Text box for Net Income
    txBox = slide1.shapes.add_textbox(left=Inches(8.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Net Income'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)



    prs.save('Earnings.pptx')
    os.startfile("Earnings.pptx")

if __name__ == '__main__':
    
    create_pptx_report4()