from cProfile import label
import numpy as np
import requests
import json
import matplotlib as mpl
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR
from numpy import shape
from pptx.util import Pt
import os

# type: ignore
# pyright: reportMissingModuleSource=false
# pyright: reportMissingImports = false

repo_rev, repo_year, ebitda, ebit, capex, fcf = [], [], [], [], [], []

prs = Presentation()
#################################################################################################################################
#################################################################################################################################

def revenue():
    
    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)

    n = len(response.json()['timeSeries']['annualTotalRevenue'])

    for i in range(n):
        repo_year.append(int(response.json()['timeSeries']['annualTotalRevenue'][i]['asOfDate'][:4]))
        repo_rev.append(float(response.json()['timeSeries']['annualTotalRevenue'][i]['reportedValue']['fmt'][:-1]))

    print(repo_year)
    print(repo_rev)

    const = 1/(repo_year[-1]-repo_year[0])
    cagr = (pow((repo_rev[-1]/repo_rev[0]), const) - 1)*100

    # CAGR calculation
    print("CAGR is :", cagr)

    # EBITDA Margin

    # json output
    with open("output.json", "w") as f:
        for i in range(n):
            f.write('"Year" : ')
            json.dump(
                int(response.json()['timeSeries']['annualTotalRevenue'][i]['asOfDate'][:4]), f)
            f.write(' , "Annual total Revenue" : ')
            json.dump(float(response.json()[
                    'timeSeries']['annualTotalRevenue'][i]['reportedValue']['fmt'][:-1]), f)
            f.write(' "Billions" \n\n')
        f.write('"CAGR is" : ')
        json.dump(float(cagr), f)
        f.write(' % \n')

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(repo_rev))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, repo_rev, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(repo_rev):
        ax.text(i-0.225,v+12, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('Revenue                                                                                  ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("CAGR.png")



def Ebitda():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    ebitda_per = []

    n = len(response.json()['timeSeries']['annualEbitda'])

    for i in range(n) :
        ebitda.append(float(response.json()['timeSeries']['annualEbitda'][i]['reportedValue']['fmt'][:-1]))

    for i in range(n) :
        ebitda_per.append(float(ebitda[i]/repo_rev[i])*100)

    print(ebitda)
    print(ebitda_per)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(ebitda))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ebitda, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ebitda):
        ax.text(i-0.2, v+4, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('EBITDA                                                                                    ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("EBITDA.png")



def Ebit():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    ebit_per = []

    n = len(response.json()['timeSeries']['annualOperatingIncome'])

    for i in range(n) :
        ebit.append(float(response.json()['timeSeries']['annualOperatingIncome'][i]['reportedValue']['fmt'][:-1]))
    
    for i in range(n) :
        ebit_per.append(float(ebit[i]/repo_rev[i])*100)

    print(ebit)
    print(ebit_per)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(ebit))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ebit, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ebit):
        ax.text(i-0.2, v+3.5, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('EBIT                                                                                         ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("EBIT.png")



def Capex():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-cash-flow"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)

    capex_per = []

    n = len(response.json()['cashflowStatementHistory']['cashflowStatements'])

    for i in range(n):
        capex.append(abs(float(response.json()['cashflowStatementHistory']['cashflowStatements'][i]['capitalExpenditures']['fmt'][:-1])))
    
    capex.reverse()
    for i in range(n) :
        capex_per.append(float(capex[i]/repo_rev[i])*100)

    print(capex)
    print(capex_per)

    '''with open("output1.json", "w") as f:
        for i in range(n):
            json.dump(abs(float(response.json()['cashflowStatementHistory']['cashflowStatements'][i]['capitalExpenditures']['fmt'][:-1])), f)'''
    
    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(capex))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, capex, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(capex):
        ax.text(i-0.2, v+0.5, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('CAPEX                                                                                     ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("CAPEX.png")



def Fcf():

    fcf_per, fcf_conversion = [], []

    n = len(repo_year)

    for i in range(n):
        fcf.append(ebitda[i]-capex[i])
        fcf_conversion.append(float(fcf[i]/ebitda[i]))
    
    for i in range(n) :
        fcf_per.append(float(fcf[i]/repo_rev[i])*100)
        fcf_conversion.append(float(fcf[i]/ebitda[i]))

    print(fcf)
    print(fcf_per)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(fcf))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, fcf, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(fcf):
        ax.text(i-0.2, v+4, str(v))

    ax.set_xticks(ind, repo_year)
    #ax.set_title('FCF                                                                                        ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("FCF.png")



def create_pptx_report1():

    revenue()
    Ebitda()
    Ebit()
    Capex()
    Fcf()

    #prs = Presentation()
    blank_slide_layout1 = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout1)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    connector = slide1.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.1), Inches(1.9), Inches(1.1))
    color_line = connector.line
    #color_line.solid()
    color_line.color.rgb = RGBColor(250, 98, 28)
    
    #Text box for Marginal Outputs
    txBox = slide1.shapes.add_textbox(left=Inches(0.64), top=Inches(0.41), width=Inches(11.5), height=Inches(0.6))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Marginal Outputs'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(32)
    font.bold = False
    font.italic = None 
    font.color.rgb = RGBColor(19, 46, 87)

    pic = slide1.shapes.add_picture('CAGR.png', Inches(1.5), Inches(1.6), width=Inches(6), height=Inches(3.5))
    pic = slide1.shapes.add_picture('EBITDA.png', Inches(8.5), Inches(1.6), width=Inches(6), height=Inches(3.5))
    pic = slide1.shapes.add_picture('EBIT.png', Inches(1.5), Inches(5.3), width=Inches(6), height=Inches(3.5))
    pic = slide1.shapes.add_picture('CAPEX.png', Inches(8.5), Inches(5.3), width=Inches(6), height=Inches(3.5))
    #pic = slide1.shapes.add_picture('FCF.png', Inches(8.5), Inches(5.3), width=Inches(6), height=Inches(3.5))

    #Text box for revenue
    txBox = slide1.shapes.add_textbox(left=Inches(1.5), top=Inches(1.6), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Revenue'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)
    

    #Text box for EBITDA
    txBox = slide1.shapes.add_textbox(left=Inches(8.5), top=Inches(1.6), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EBITDA'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    #Text box for EBIT
    txBox = slide1.shapes.add_textbox(left=Inches(1.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EBIT'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    #Text box for CAPEX
    txBox = slide1.shapes.add_textbox(left=Inches(8.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'CAPEX'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)

    #Text box for FCF
    '''txBox = slide1.shapes.add_textbox(left=Inches(8.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'FCF'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)'''


    prs.save('marginal_outputs.pptx')
    os.startfile("marginal_outputs.pptx")


if __name__ == '__main__':

    create_pptx_report1()
    