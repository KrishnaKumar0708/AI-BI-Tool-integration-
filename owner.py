from cProfile import label
import numpy as np
import requests
import json
import matplotlib as mpl
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR
from numpy import shape
from pptx.util import Pt
import os

# type: ignore
# pyright: reportMissingModuleSource=false
# pyright: reportMissingImports = false


ow_st = []

prs = Presentation()

#################################################################################################################################
#################################################################################################################################


def Owner_Struct():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-holders"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)
    #n = len(response.json()['timeSeries']['annualBasicEPS'])
    Labels = ["Insider's",'Institutional\nInvestors', 'Other']
    #Colors = ['#203864','#2F5597','#8FAADC']
    Colors = ['#FFFFFF','#132E57','#8FAADC']
    newest = []
    ow_st.append(float(response.json()['majorHoldersBreakdown']['insidersPercentHeld']['fmt'][:-1]))
    ow_st.append(float(response.json()['majorHoldersBreakdown']['institutionsPercentHeld']['fmt'][:-1]))
    ow_st.append(round((100 - ow_st[0] - ow_st[1]), 2))

    print(ow_st)

    '''with open("output.json", "w") as f:
        json.dump(response.json()['majorHoldersBreakdown']['insidersPercentHeld'], f)
        json.dump(response.json()['majorHoldersBreakdown']['institutionsPercentHeld'], f)'''
    
    n = len(ow_st)
    for i in range(n):
        newest.append(Labels[i]+ '\n'+ str(ow_st[i]) + '%')
    fig, ax = plt.subplots()
    ax.pie(ow_st, labels = newest, labeldistance = 1.25, colors = Colors, textprops = dict(ha="center"), wedgeprops = {"linewidth": 2, "edgecolor": "white"})
    plt.savefig("Ownership_Structure.png")
    


def create_pptx_report3():
    
    Owner_Struct()
    
    #prs = Presentation()
    blank_slide_layout3 = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(blank_slide_layout3)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    connector = slide3.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.1), Inches(1.9), Inches(1.1))
    color_line = connector.line
    #color_line.solid()
    color_line.color.rgb = RGBColor(250, 98, 28)
    
    #Text box for Ownership Structure
    txBox = slide3.shapes.add_textbox(left=Inches(0.64), top=Inches(0.41), width=Inches(11.5), height=Inches(0.6))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Ownership Structure'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(32)
    font.bold = False
    font.italic = None 
    font.color.rgb = RGBColor(19, 46, 87)
    

    #slide.AddConnector(msoConnectorStraight, Inches(0), Inches(0.5), Inches(0.5), Inches(0.5))
    pic = slide3.shapes.add_picture('Ownership_Structure.png', Inches(4.5), Inches(2.9))
    

    #Text box for Ownership Structure
    txBox = slide3.shapes.add_textbox(left=Inches(4.5), top=Inches(2.5), width=Inches(7), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Ownership Structure'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    prs.save('ownership_structure.pptx')
    os.startfile("ownership_structure.pptx")

#################################################################################################################################
#################################################################################################################################

if __name__ == '__main__':
    
    create_pptx_report3()