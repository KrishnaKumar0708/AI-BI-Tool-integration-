from cProfile import label
import numpy as np
import requests
import json
import matplotlib as mpl
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR
from numpy import shape
from pptx.util import Pt
import os

# type: ignore
# pyright: reportMissingModuleSource=false
# pyright: reportMissingImports = false

repo_rev, repo_year, ebitda, ebit, capex, fcf = [], [], [], [], [], []

ev, ev_sales, ev_ebitda, ev_ebit, ev_fcf = [], [], [], [], []

stock_price, earnings, p_e = [], [], []

ow_st = []

prs = Presentation()

#################################################################################################################################
#################################################################################################################################

def revenue():
    
    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)

    n = len(response.json()['timeSeries']['annualTotalRevenue'])

    for i in range(n):
        repo_year.append(int(response.json()['timeSeries']['annualTotalRevenue'][i]['asOfDate'][:4]))
        repo_rev.append(float(response.json()['timeSeries']['annualTotalRevenue'][i]['reportedValue']['fmt'][:-1]))

    


def Ebitda():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    ebitda_per = []

    n = len(response.json()['timeSeries']['annualEbitda'])

    for i in range(n) :
        ebitda.append(float(response.json()['timeSeries']['annualEbitda'][i]['reportedValue']['fmt'][:-1]))

    for i in range(n) :
        ebitda_per.append(float(ebitda[i]/repo_rev[i])*100)

    


def Ebit():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol":"AAPL","region":"US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    ebit_per = []

    n = len(response.json()['timeSeries']['annualOperatingIncome'])

    for i in range(n) :
        ebit.append(float(response.json()['timeSeries']['annualOperatingIncome'][i]['reportedValue']['fmt'][:-1]))
    
    for i in range(n) :
        ebit_per.append(float(ebit[i]/repo_rev[i])*100)

    



def Capex():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-cash-flow"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)

    capex_per = []

    n = len(response.json()['cashflowStatementHistory']['cashflowStatements'])

    for i in range(n):
        capex.append(abs(float(response.json()['cashflowStatementHistory']['cashflowStatements'][i]['capitalExpenditures']['fmt'][:-1])))
    
    for i in range(n) :
        capex_per.append(float(capex[i]/repo_rev[i])*100)

    



def Fcf():

    fcf_per = []

    n = len(repo_year)

    for i in range(n):
        fcf.append(ebitda[i]-capex[i])
    
    for i in range(n) :
        fcf_per.append(float(fcf[i]/repo_rev[i])*100)

    





#################################################################################################################################
#################################################################################################################################

def EV():
    
    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-statistics"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)

    n = len(repo_year)
    ev.append(float(response.json()['defaultKeyStatistics']['enterpriseValue']['fmt'][:-1]))
    
    print(ev)

    '''with open("output.json", "w") as f:
        for i in range(n):
            json.dump(response.json(), f)'''
            


def EV_SALES():

    n = len(repo_year)
    for i in range(n):
        ev_sales.append(round(float(ev[0]*1000/repo_rev[i]), 2))
    
    print(ev_sales)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(ev_sales))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ev_sales, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ev_sales):
        ax.text(i-0.2, v+0.5, str(v)+'x')

    ax.set_xticks(ind, repo_year)
    #ax.set_title('FCF                                                                                        ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    #ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Ev_Sales.png")



def EV_EBITDA():

    n = len(repo_year)
    for i in range(n):
        ev_ebitda.append(round(float(ev[0]*1000/ebitda[i]), 2))
    
    print(ev_ebitda)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(ev_ebitda))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ev_ebitda, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ev_ebitda):
        ax.text(i-0.2, v+1.5, str(v)+'x')

    ax.set_xticks(ind, repo_year)
    #ax.set_title('FCF                                                                                        ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    #ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Ev_Ebitda.png")



def EV_EBIT():

    n = len(repo_year)
    for i in range(n):
        ev_ebit.append(round(float(ev[0]*1000/ebit[i]), 2))
    
    print(ev_ebit)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(ev_ebit))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ev_ebit, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ev_ebit):
        ax.text(i-0.2, v+1.5, str(v)+'x')

    ax.set_xticks(ind, repo_year)
    #ax.set_title('FCF                                                                                        ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    #ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Ev_Ebit.png")



def EV_FCF():

    n = len(repo_year)
    for i in range(n):
        ev_fcf.append(round(float(ev[0]*1000/fcf[i]), 2))

    print(ev_fcf)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(ev_fcf))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.2)

    ax.bar(ind, ev_fcf, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(ev_fcf):
        ax.text(i-0.2, v+1.5, str(v)+'x')

    ax.set_xticks(ind, repo_year)
    #ax.set_title('FCF                                                                                        ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    #ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("Ev_Fcf.png")



def stock_Price():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-statistics"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)
    n = len(response.json()['price']['regularMarketPrice'])

    stock_price.append(response.json()['price']['regularMarketPrice']['raw'])
    
    print(stock_price)



def Earnings():

    url = "https://yh-finance.p.rapidapi.com/stock/v2/get-financials"

    querystring = {"symbol": "AAPL", "region": "US"}

    headers = {
        'x-rapidapi-host': "yh-finance.p.rapidapi.com",
        'x-rapidapi-key': "e7f3bf2b7cmsh06ecd7e8b764151p11c4ecjsnddfceea6a81f"
    }

    response = requests.request("GET", url, headers=headers, params=querystring)
    n = len(response.json()['timeSeries']['annualBasicEPS'])

    for i in range(n):
        earnings.append(float(response.json()['timeSeries']['annualBasicEPS'][i]['reportedValue']['raw']))
    
    print(earnings)
        


def P_E():

    n = len(earnings)
    for i in range(n):
        p_e.append(round((stock_price[0]/earnings[i]), 2))
    
    print(p_e)

    #Bar plot
    fig, ax = plt.subplots()
    width = 0.4
    ind = np.arange(len(p_e))

    ax.set_xmargin(0.2)
    ax.set_ymargin(0.1)

    ax.bar(ind, p_e, width, color='#132E57')
    plt.tick_params(
        axis='both',          # changes apply to the x-axis
        which='both',      # both major and minor ticks are affected
        bottom=False,      # ticks along the bottom edge are off
        top=False,
        left=False,
        right=False,        # ticks along the top edge are off
        labelbottom=True)

    # ax.set_xticks(repo_year,repo_year)
    for i, v in enumerate(p_e):
        ax.text(i-0.2, v+1.5, str(v)+'x')

    ax.set_xticks(ind, repo_year)
    #ax.set_title('FCF                                                                                        ', color='white',
                #fontweight='bold', loc='left', x=-0.12, y=1.05, backgroundcolor='#132E57', fontsize=12)
    # ax.set_xlabel("Years")
    #ax.set_ylabel("(Billions of USD)")

    for spine in ax.spines:
        ax.spines[spine].set_visible(False)
    plt.gca().spines['bottom'].set_visible(True)
    #plt.gca().spines['top'].set_visible(True)
    plt.gca().spines['bottom'].set_color('#BDD1F1')

    plt.savefig("P_E.png")



def create_pptx_report2():

    revenue()
    Ebitda()
    Ebit()
    Capex()
    Fcf()
    EV()
    EV_SALES()
    EV_EBITDA()
    EV_EBIT()
    EV_FCF()
    stock_Price()
    Earnings()
    P_E()
    
    #prs = Presentation()
    blank_slide_layout2 = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(blank_slide_layout2)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    connector = slide2.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.1), Inches(1.9), Inches(1.1))
    color_line = connector.line
    #color_line.solid()
    color_line.color.rgb = RGBColor(250, 98, 28)
    
    #Text box for Valuation Outputs
    txBox = slide2.shapes.add_textbox(left=Inches(0.64), top=Inches(0.41), width=Inches(11.5), height=Inches(0.6))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'Valuation Outputs'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(32)
    font.bold = False
    font.italic = None 
    font.color.rgb = RGBColor(19, 46, 87)

    pic = slide2.shapes.add_picture('Ev_Sales.png', Inches(1.5), Inches(1.6), width=Inches(6), height=Inches(3.5))
    pic = slide2.shapes.add_picture('Ev_Ebitda.png', Inches(8.5), Inches(1.6), width=Inches(6), height=Inches(3.5))
    pic = slide2.shapes.add_picture('Ev_Ebit.png', Inches(1.5), Inches(5.3), width=Inches(6), height=Inches(3.5))
    #pic = slide2.shapes.add_picture('Ev_Fcf.png', Inches(8.5), Inches(5.3), width=Inches(6), height=Inches(3.5))
    pic = slide2.shapes.add_picture('P_E.png', Inches(8.5), Inches(5.3), width=Inches(6), height=Inches(3.5))

    #Text box for revenue
    txBox = slide2.shapes.add_textbox(left=Inches(1.5), top=Inches(1.6), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EV/SALES'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)
    

    #Text box for EBITDA
    txBox = slide2.shapes.add_textbox(left=Inches(8.5), top=Inches(1.6), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EV/EBITDA'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    #Text box for EBIT
    txBox = slide2.shapes.add_textbox(left=Inches(1.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EV/EBIT'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    '''#Text box for FCF
    txBox = slide2.shapes.add_textbox(left=Inches(8.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'EV/FCF'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)'''


    #Text box for P_E
    txBox = slide2.shapes.add_textbox(left=Inches(8.5), top=Inches(5.3), width=Inches(6), height=Inches(0.4))
    tx = txBox.text_frame
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = 'P/E'
    font = run.font 
    font.name = 'Open Sans Light'
    font.size = Pt(16)
    font.bold = True
    font.italic = None 
    font.color.rgb = RGBColor(255, 255, 255)
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(19, 46, 87)


    prs.save('valuation_outputs.pptx')
    os.startfile("valuation_outputs.pptx")

#################################################################################################################################
#################################################################################################################################



#################################################################################################################################
#################################################################################################################################

if __name__ == '__main__':
    
    create_pptx_report2()
    